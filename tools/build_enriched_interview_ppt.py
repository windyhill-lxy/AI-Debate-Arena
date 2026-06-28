from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMP_DIR = Path(r"D:\临时文件")
TEMPLATE = TEMP_DIR / "模板版.pptx"
OUTPUT = TEMP_DIR / "模板版-补全可编辑版.pptx"
ASSET_DIR = TEMP_DIR / "AI辩论场PPT补全素材"

FLOW_SPECS = {
    "overall": (
        "从真实训练到可复盘系统",
        ["训练痛点", "辩题资料", "论据入库", "赛制推进", "评分复盘"],
    ),
    "architecture": (
        "项目运行架构",
        ["React前端", "WebSocket", "FastAPI", "LangGraph", "模型与RAG", "导出复盘"],
    ),
    "rag": (
        "防幻觉论据链路",
        ["资料检索", "来源绑定", "引用校验", "事实核查", "回环重写", "裁判扣分"],
    ),
    "interview": (
        "访谈反馈进入迭代",
        ["真实辩手体验", "记录痛点", "校准赛制", "调整边界", "回到测试"],
    ),
    "schedule": (
        "57步正式赛程",
        ["开场准备", "立论", "驳立论", "质辩", "自由辩", "总结裁判"],
    ),
    "delivery": (
        "可交付验证路径",
        ["一键启动", "运行设置", "自动化测试", "联机演示", "报告导出"],
    ),
}


COLORS = {
    "ink": "1C2233",
    "muted": "5B6478",
    "blue": "243A5E",
    "teal": "2D7D73",
    "orange": "C46A32",
    "cream": "F7F2E8",
    "white": "FFFFFF",
    "line": "C9D2E3",
    "pale_blue": "EAF2FF",
    "pale_teal": "E9F6F3",
    "pale_orange": "FFF0E7",
}


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def set_run_font(run, size: float, color: str = "ink", bold: bool | None = None) -> None:
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.color.rgb = rgb(COLORS[color])
    if bold is not None:
        run.font.bold = bold


def set_shape_text_style(shape, size: float, color: str = "ink", bold: bool | None = None) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.line_spacing = 0.9
        paragraph.space_after = Pt(1)
        for run in paragraph.runs:
            set_run_font(run, size, color, bold)


def text_box(slide, x, y, w, h, title: str, body: str, *, fill="white", title_color="blue") -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[fill])
    shape.fill.transparency = 4
    shape.line.color.rgb = rgb(COLORS["line"])
    shape.line.width = Pt(0.9)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run_font(r, 12.8, title_color, True)
    p.space_after = Pt(4)
    p2 = tf.add_paragraph()
    p2.text = body
    p2.line_spacing = 1.05
    p2.space_after = Pt(0)
    for run in p2.runs:
        set_run_font(run, 10.6, "ink", False)


def plain_label(slide, x, y, w, h, text: str, *, size=9.2, color="muted", bold=False) -> None:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        set_run_font(run, size, color, bold)


def screenshot_box(slide, x, y, w, h, instruction: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS["white"])
    shape.fill.transparency = 8
    shape.line.color.rgb = rgb(COLORS["orange"])
    shape.line.width = Pt(1.3)
    shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = "截图预留"
    for run in p.runs:
        set_run_font(run, 12.4, "orange", True)
    p2 = tf.add_paragraph()
    p2.text = instruction
    p2.line_spacing = 1.0
    for run in p2.runs:
        set_run_font(run, 9.2, "ink", False)


def slide_footer(slide, n: int) -> None:
    plain_label(
        slide,
        0.6,
        7.17,
        11.6,
        0.18,
        f"补充内容依据项目文件、赛程配置、代码模块与何彩嫦访谈整理｜第 {n:02d} 页",
        size=7.6,
        color="muted",
    )


def split_label(text: str, max_len: int = 7) -> list[str]:
    if "·" in text or " " in text:
        parts = re.split(r"[ ·]", text)
        lines: list[str] = []
        current = ""
        for part in parts:
            if not part:
                continue
            if len(current + part) <= max_len:
                current += part
            else:
                if current:
                    lines.append(current)
                current = part
        if current:
            lines.append(current)
        return lines[:3]
    return [text[i : i + max_len] for i in range(0, len(text), max_len)][:3]


def create_flow_svg(path: Path, title: str, nodes: list[str], *, colors: list[str] | None = None) -> None:
    width, height = 1120, 310
    margin = 54
    gap = 24
    node_w = (width - margin * 2 - gap * (len(nodes) - 1)) / len(nodes)
    node_h = 92
    y = 126
    colors = colors or ["#EAF2FF", "#E9F6F3", "#FFF0E7", "#F7F2E8", "#EEF0F8"]
    parts = [
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}" viewBox="0 0 {width} {height}">',
        '<rect width="1120" height="310" rx="28" fill="#FBFCFE"/>',
        '<text x="54" y="58" font-family="Microsoft YaHei, SimHei, sans-serif" font-size="30" font-weight="700" fill="#1C2233">'
        + title
        + "</text>",
    ]
    for i, node in enumerate(nodes):
        x = margin + i * (node_w + gap)
        if i < len(nodes) - 1:
            ax1 = x + node_w + 4
            ax2 = x + node_w + gap - 8
            ay = y + node_h / 2
            parts.append(
                f'<line x1="{ax1:.1f}" y1="{ay:.1f}" x2="{ax2:.1f}" y2="{ay:.1f}" '
                'stroke="#6E7B91" stroke-width="4" stroke-linecap="round"/>'
            )
            parts.append(
                f'<path d="M {ax2:.1f} {ay:.1f} l -12 -8 v 16 z" fill="#6E7B91"/>'
            )
        parts.append(
            f'<rect x="{x:.1f}" y="{y}" width="{node_w:.1f}" height="{node_h}" rx="18" '
            f'fill="{colors[i % len(colors)]}" stroke="#9DA9BC" stroke-width="2"/>'
        )
        lines = split_label(node)
        start_y = y + 38 - (len(lines) - 1) * 13
        for j, line in enumerate(lines):
            parts.append(
                f'<text x="{x + node_w / 2:.1f}" y="{start_y + j * 27:.1f}" '
                'text-anchor="middle" font-family="Microsoft YaHei, SimHei, sans-serif" '
                'font-size="22" font-weight="700" fill="#1C2233">'
                + line
                + "</text>"
            )
    parts.append("</svg>")
    path.write_text("\n".join(parts), encoding="utf-8")


def prepare_assets() -> dict[str, tuple[str, list[str]]]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    out: dict[str, tuple[str, list[str]]] = {}
    for key, (title, nodes) in FLOW_SPECS.items():
        svg_path = ASSET_DIR / f"{key}.svg"
        create_flow_svg(svg_path, title, nodes)
        out[key] = (title, nodes)
    return out


def add_flow(slide, flow_spec: tuple[str, list[str]], x, y, w, h) -> None:
    title, nodes = flow_spec
    title_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.24))
    title_box.text_frame.text = title
    set_shape_text_style(title_box, 9.8, "blue", True)

    n = len(nodes)
    gap = 0.1 if n >= 6 else 0.16
    arrow_w = 0.16 if n >= 6 else 0.18
    node_w = (w - gap * (n - 1) - arrow_w * (n - 1)) / n
    node_h = max(0.45, h - 0.38)
    node_y = y + 0.34
    palette = ["pale_blue", "pale_teal", "pale_orange", "cream", "white", "pale_blue"]

    # Add arrows first so node shapes and labels remain visually on top.
    for i in range(n - 1):
        arrow_x = x + (i + 1) * node_w + i * (gap + arrow_w) + gap * 0.55
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(arrow_x),
            Inches(node_y + node_h * 0.33),
            Inches(arrow_w),
            Inches(node_h * 0.34),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = rgb(COLORS["muted"])
        arrow.line.color.rgb = rgb(COLORS["muted"])

    for i, node in enumerate(nodes):
        node_x = x + i * (node_w + gap + arrow_w)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(node_x),
            Inches(node_y),
            Inches(node_w),
            Inches(node_h),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = rgb(COLORS[palette[i % len(palette)]])
        box.fill.transparency = 3
        box.line.color.rgb = rgb(COLORS["line"])
        box.line.width = Pt(0.8)
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = Inches(0.03)
        tf.margin_right = Inches(0.03)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        p = tf.paragraphs[0]
        p.text = node
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            set_run_font(run, 8.4 if n >= 6 else 9.2, "ink", True)


def reposition_original(slide, x, y, w, h, font_size: float, *, skip_first_title=False) -> None:
    text_shapes = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    ]
    if not text_shapes:
        return
    target_shapes = text_shapes[1:] if skip_first_title and len(text_shapes) > 1 else text_shapes
    for shape in target_shapes:
        shape.left = Inches(x)
        shape.top = Inches(y)
        shape.width = Inches(w)
        shape.height = Inches(h)
        set_shape_text_style(shape, font_size)


def normalize_title_slide(slide, font_size=34) -> None:
    text_shapes = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    ]
    if not text_shapes:
        return
    shape = text_shapes[0]
    shape.left = Inches(0.7)
    shape.top = Inches(0.45)
    shape.width = Inches(11.9)
    shape.height = Inches(0.75)
    set_shape_text_style(shape, font_size, "blue", True)


SLIDES = [
    {
        "flow": "overall",
        "title": "补充说明｜为什么这个作品适合智能体专项赛",
        "body": "我们做的不是把辩题丢给大模型再让它连续输出，而是把辩论拆成可执行的流程：谁在什么时候发言、能看到哪些资料、引用从哪里来、发言后怎样核查，全部落在程序状态里。这样处理以后，AI 辩手不只是会说话，它要按赛制、按角色、按证据行动，评委也能看到过程记录。",
        "shot": "请粘贴作品首页或辩论室总览截图，最好能同时看到项目名称、辩题入口和三栏辩论室。",
        "layout": "cover",
    },
    {
        "flow": "overall",
        "title": "补充说明｜立项不是为了做一个聊天窗口",
        "body": "访谈里最直接的反馈是：辩论训练主要靠赛前模拟和赛后老师指出问题，同学真正缺的是稳定练习环境、现场抓关键词能力、队友之间的配合，以及有来源的论据材料。项目立项时，我们把这些痛点转成系统目标：凑不齐人时 AI 能补位，发言结束后能复盘，论据不能凭空编，用户也能加入真实赛程练习。",
        "shot": "请粘贴辩论训练、社团活动或访谈整理文件截图，用来证明需求来自真实使用场景。",
        "layout": "section",
    },
    {
        "flow": "interview",
        "title": "补充说明｜访谈把“练不起来”的问题说清楚了",
        "body": "何彩嫦同学提到，比赛前通常会用同一个辩题做两到三次模拟，再根据老师指出的问题改下一场。二辩最难的是现场抓住对方关键词，把反驳顺成能说出口的话；自由辩虽然是即时攻防，但队友能互相接力，不能只让一个人全程承担压力。这些话让我们确定系统必须支持正式赛程、队内讨论和人类随时接入。",
        "shot": "请粘贴访谈 Word 文档或访谈现场照片，位置建议放录音转文字的关键段落。",
        "layout": "dense",
    },
    {
        "flow": "rag",
        "title": "补充说明｜同质化问题的本质是没有边界",
        "body": "调研后发现，许多 AI 辩论项目只是把双方身份、辩题和历史消息一次性塞进 prompt。这样看起来能对话，但 AI 可能知道不该知道的信息，也没有稳定的证据来源，更没有赛程状态。我们在后端用 ai_context_manager 管理每个角色能看的内容，用 RAG 和引用校验限制事实表达，避免它变成“口才很好的聊天机器人”。",
        "shot": "请粘贴 GitHub 调研截图、竞品页面截图，或旧版 AI 输出中出现无来源材料的例子。",
        "layout": "dense",
    },
    {
        "flow": "rag",
        "title": "补充说明｜论据库、多模态和用户参与要一起解决",
        "body": "访谈中对论据的要求很具体：最好是大众能听懂的真实数据和案例，来源要明确，过度生僻的科学解释反而浪费发言时间；AI 还容易滥用比喻，论证力度不足。对应到系统里，我们把论据入库、引用详情、语音录入、TTS 朗读、摄像头表现评价和用户席位都纳入同一个辩论流程，而不是分散成几个演示功能。",
        "shot": "请粘贴论据库、引用详情面板、语音录入按钮或摄像头训练窗口截图。",
        "layout": "dense",
    },
    {
        "flow": "architecture",
        "title": "补充说明｜基本架构和运行方式",
        "body": "项目后端是 FastAPI，负责房间、赛程、WebSocket、导出和管理接口；前端是 React + Vite，负责首页、辩论室、回放、管理页和联机加入流程；工作流由 LangGraph 风格的 DebateGraph 推进。MongoDB、Redis 可选，没有数据库时会自动降级到内存模式，方便在学校机房或演示电脑上快速启动。",
        "shot": "请粘贴首页创建房间界面，或应用程序一键启动后的浏览器首页截图。",
        "layout": "section",
    },
    {
        "flow": "architecture",
        "title": "补充说明｜系统基本介绍",
        "body": "用户可以选择 AI 自主辩论、加入正方、加入反方和多人联机辩论。进入房间后，左栏管理资料、赛程和视角，中央显示正式发言和流式输出，右栏展示队内讨论、AI 策略和当前回合信息。辩论结束后，系统能导出 Markdown 或 PDF 复盘报告，记录发言、引用、评分理由和裁判总结。",
        "shot": "请粘贴辩论室主界面截图，建议能看到左栏、中央发言区、右栏队内讨论三个区域。",
        "layout": "blank",
    },
    {
        "flow": "overall",
        "title": "补充说明｜三个创新点互相支撑",
        "body": "“更真实”解决赛制和体验，“更严谨”解决证据和角色边界，“更实用”解决真实部署和训练价值。它们不是三个分开的卖点：如果没有正式流程，防幻觉只是在聊天里做检查；如果没有论据和上下文权限，多智能体容易串台；如果用户不能加入、不能联机、不能导出报告，系统也很难成为可用的训练工具。",
        "shot": "请粘贴项目流程图弹窗、三栏辩论室或能力总览页截图，用来承接三个创新点。",
        "layout": "section",
    },
    {
        "flow": "schedule",
        "title": "补充说明｜“真实”体现在赛程细节里",
        "body": "formal_4v4.yaml 里配置了 57 个赛程段，从开场、开场论据入库、双方一辩任务分配，到立论、驳立论、质辩、自由辩、总结陈词和裁判终局。访谈里提到自由辩是一方落座后另一方立即计时、每个人都要有发言，这正好对应系统中的赛程推进、用户席位判断和队内讨论补位。",
        "shot": "请粘贴赛程进度条、流程图高亮节点、队内讨论或自由辩页面截图。",
        "layout": "wide",
    },
    {
        "flow": "rag",
        "title": "补充说明｜“严谨”不是口号，是每轮发言前后的约束",
        "body": "后端每轮发言会经过检索、策略、方向判断、反思、生成、事实核查、发布、评分和推进。公开发言中的引用会进入 sanitize_citations 检查，未入库的引用标记会被移除；涉及数据、法规、研究表明等高风险表述时，fact_check 会要求回环重写或在评分中扣幻觉风险分。访谈里对“数据来源”和“比喻不能代替论证”的要求，正好落在这一层。",
        "shot": "请粘贴引用详情面板、发言中的引用标记、论据来源摘录或裁判扣分说明截图。",
        "layout": "wide",
    },
    {
        "flow": "delivery",
        "title": "补充说明｜“实用”体现在容错和可分发",
        "body": "项目不仅追求功能跑通，还处理了演示时最常见的问题：没有摄像头、网络不通、API Key 未填写、TTS/ASR 不可用、联机用户未就绪等情况会显示明确提示，而不是直接崩溃。应用程序目录里保留了一键启动和打包脚本，运行设置面板可以配置模型密钥，多人联机也支持主持控制和访客加入。",
        "shot": "请粘贴运行设置面板、错误弹窗、联机等待页、主持控制或一辩改稿功能截图。",
        "layout": "wide",
    },
    {
        "flow": "delivery",
        "title": "补充说明｜开发历程按“能跑、可信、好用”推进",
        "body": "第一阶段先跑通基本对话和房间状态，验证 AI 轮流发言；第二阶段把赛程拆成节点，引入论据检索和 LangGraph 式循环；第三阶段补上语音朗读、语音录入、队内讨论、摄像头表现评价和联机能力。后期重点转向工程可靠性，包括 39 个后端测试文件、前端构建、E2E 脚本、便携启动和导出报告。",
        "shot": "请粘贴测试命令通过截图、项目目录、启动脚本、开发记录或 Git 提交记录截图。",
        "layout": "dev",
    },
    {
        "flow": "interview",
        "title": "补充说明｜访谈反馈怎样改变系统",
        "body": "访谈没有只问“你觉得好不好”，而是围绕辩位、赛制、论据、自由辩和评委复盘继续追问。她指出一辩需要撑满三分钟、要有三个论点和具体数据；二辩要现场抓关键词；自由辩不能只有一个人发言；评委理想状态应该指出逻辑漏洞而不是只说客套话。这些反馈让系统更重视赛程节奏、论据质量和复盘报告。",
        "shot": "请粘贴访谈现场照片、访谈转写重点标注，或把反馈整理成表格后的截图。",
        "layout": "section",
    },
    {
        "flow": "overall",
        "title": "补充说明｜尾声：把一次作品做成可继续训练的工具",
        "body": "最终交付的不是单一演示页面，而是一套可运行、可测试、可分发、可回放的辩论训练系统。它保留了学生项目的探索性：从一开始的“AI 能不能辩论”，逐步变成“怎样让 AI 在清晰规则下帮助人练辩论”。后续如果继续迭代，重点会放在更多赛制模板、更稳定的语音实时辩论和更多辩论社回访记录。",
        "shot": "请粘贴最终演示画面、PDF 复盘报告、回放页或答辩现场照片。",
        "layout": "section",
    },
    {
        "flow": "delivery",
        "title": "补充说明｜最后一页可放截图清单和演示路径",
        "body": "答辩时建议按这条线讲：先说真实训练痛点，再说为什么普通聊天机器人不够；接着展示创建房间、论据入库、57 步赛程、RAG 引用、防幻觉核查、用户加入和导出报告；最后用访谈反馈说明系统不是闭门造车。截图可以按“首页、辩论室、流程图、引用详情、联机、导出报告”六类准备。",
        "shot": "请粘贴作品二维码、最终运行截图合集，或放一张最能代表项目完成度的总览图。",
        "layout": "blank",
    },
]


def add_content_for_slide(slide, idx: int, spec: dict, flows: dict[str, Path]) -> None:
    layout = spec["layout"]
    if layout in {"section", "blank"}:
        normalize_title_slide(slide)
        text_box(slide, 0.75, 1.55, 6.45, 2.75, spec["title"], spec["body"], fill="cream")
        add_flow(slide, flows[spec["flow"]], 7.55, 1.55, 4.9, 1.36)
        screenshot_box(slide, 7.55, 3.15, 4.9, 2.5, spec["shot"])
    elif layout == "cover":
        text_box(slide, 0.72, 3.22, 4.9, 2.45, spec["title"], spec["body"], fill="cream")
        add_flow(slide, flows[spec["flow"]], 6.05, 2.7, 5.9, 1.63)
        screenshot_box(slide, 6.05, 4.55, 5.9, 1.45, spec["shot"])
    elif layout == "dense":
        reposition_original(slide, 0.42, 0.42, 5.75, 6.42, 7.8)
        text_box(slide, 6.35, 0.55, 3.55, 3.45, spec["title"], spec["body"], fill="pale_blue")
        screenshot_box(slide, 10.15, 0.55, 2.75, 2.4, spec["shot"])
        add_flow(slide, flows[spec["flow"]], 6.35, 4.35, 6.55, 1.82)
    elif layout == "wide":
        title_shapes = [
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        ]
        if title_shapes:
            title_shapes[0].left = Inches(0.75)
            title_shapes[0].top = Inches(0.35)
            title_shapes[0].width = Inches(11.7)
            title_shapes[0].height = Inches(1.1)
            set_shape_text_style(title_shapes[0], 26, "blue", True)
        if len(title_shapes) > 1:
            body = title_shapes[1]
            body.left = Inches(0.75)
            body.top = Inches(1.75)
            body.width = Inches(5.25)
            body.height = Inches(3.9)
            set_shape_text_style(body, 11.2)
        text_box(slide, 6.28, 1.58, 3.65, 3.32, spec["title"], spec["body"], fill="pale_teal")
        screenshot_box(slide, 10.18, 1.58, 2.7, 2.35, spec["shot"])
        add_flow(slide, flows[spec["flow"]], 6.28, 5.12, 6.6, 1.52)
    elif layout == "dev":
        reposition_original(slide, 0.7, 0.6, 4.7, 4.35, 13.2)
        text_box(slide, 5.65, 0.72, 4.05, 3.25, spec["title"], spec["body"], fill="pale_orange")
        screenshot_box(slide, 9.98, 0.72, 2.85, 2.65, spec["shot"])
        add_flow(slide, flows[spec["flow"]], 5.65, 4.35, 7.18, 1.62)
    else:
        text_box(slide, 0.75, 1.55, 6.45, 2.75, spec["title"], spec["body"], fill="cream")
        add_flow(slide, flows[spec["flow"]], 7.55, 1.55, 4.9, 1.36)
        screenshot_box(slide, 7.55, 3.15, 4.9, 2.5, spec["shot"])
    slide_footer(slide, idx)


def add_blank_slide_7_title(slide) -> None:
    if any(getattr(shape, "has_text_frame", False) and shape.text.strip() for shape in slide.shapes):
        return
    title = slide.shapes.add_textbox(Inches(0.74), Inches(0.55), Inches(11.6), Inches(0.75))
    tf = title.text_frame
    tf.text = "系统基本功能"
    set_shape_text_style(title, 34, "blue", True)


def add_blank_slide_15_title(slide) -> None:
    for shape in list(slide.shapes):
        if shape.shape_type == 13:
            shape.element.getparent().remove(shape.element)
    if any(getattr(shape, "has_text_frame", False) and shape.text.strip() for shape in slide.shapes):
        return
    title = slide.shapes.add_textbox(Inches(0.74), Inches(0.55), Inches(11.6), Inches(0.75))
    tf = title.text_frame
    tf.text = "答辩演示路径与截图清单"
    set_shape_text_style(title, 32, "blue", True)


def main() -> None:
    flows = prepare_assets()
    prs = Presentation(str(TEMPLATE))
    if len(prs.slides) != 15:
        raise RuntimeError(f"模板页数不是 15 页，当前为 {len(prs.slides)} 页")

    add_blank_slide_7_title(prs.slides[6])
    add_blank_slide_15_title(prs.slides[14])

    for idx, (slide, spec) in enumerate(zip(prs.slides, SLIDES, strict=True), start=1):
        add_content_for_slide(slide, idx, spec, flows)

    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    main()
