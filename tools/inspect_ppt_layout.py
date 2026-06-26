from pathlib import Path
import sys

from pptx import Presentation


EMU_PER_INCH = 914400


def inches(value: int) -> float:
    return round(value / EMU_PER_INCH, 2)


def main() -> int:
    if len(sys.argv) != 2:
        print("Usage: python tools/inspect_ppt_layout.py <pptx>", file=sys.stderr)
        return 2
    prs = Presentation(sys.argv[1])
    print(f"slides={len(prs.slides)} size={inches(prs.slide_width)}x{inches(prs.slide_height)}")
    for idx, slide in enumerate(prs.slides, 1):
        print(f"\nSLIDE {idx}")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text = shape.text.replace("\n", " / ")
                print(
                    f"  x={inches(shape.left)} y={inches(shape.top)} "
                    f"w={inches(shape.width)} h={inches(shape.height)} "
                    f"text={text[:160]!r}"
                )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
