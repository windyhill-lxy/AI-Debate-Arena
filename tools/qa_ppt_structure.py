from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

EMU_PER_INCH = 914400


def inch(value: int) -> float:
    return value / EMU_PER_INCH


def main() -> int:
    if len(sys.argv) != 2:
        print("Usage: python tools/qa_ppt_structure.py <pptx>", file=sys.stderr)
        return 2
    path = Path(sys.argv[1])
    prs = Presentation(str(path))
    slide_w = inch(prs.slide_width)
    slide_h = inch(prs.slide_height)
    issues: list[str] = []
    summary: list[str] = []
    flow_titles = {
        "从真实训练到可复盘系统",
        "项目运行架构",
        "防幻觉论据链路",
        "访谈反馈进入迭代",
        "57步正式赛程",
        "可交付验证路径",
    }
    for idx, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        text_shapes = 0
        picture_count = 0
        for shape in slide.shapes:
            left = inch(shape.left)
            top = inch(shape.top)
            right = inch(shape.left + shape.width)
            bottom = inch(shape.top + shape.height)
            if left < -0.01 or top < -0.01 or right > slide_w + 0.01 or bottom > slide_h + 0.01:
                issues.append(
                    f"slide {idx}: shape out of bounds "
                    f"({left:.2f},{top:.2f},{right:.2f},{bottom:.2f})"
                )
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    text_shapes += 1
                    texts.append(text)
            if shape.shape_type == 13:
                picture_count += 1
        full_text = "\n".join(texts)
        has_addition = "补充说明｜" in full_text
        has_screenshot = "截图预留" in full_text
        has_flow = any(title in full_text for title in flow_titles)
        if not has_addition:
            issues.append(f"slide {idx}: missing added explanation")
        if not has_screenshot:
            issues.append(f"slide {idx}: missing screenshot placeholder")
        if not has_flow:
            issues.append(f"slide {idx}: missing flow diagram text")
        if picture_count:
            issues.append(f"slide {idx}: contains {picture_count} picture object(s); expected editable text/shapes only")
        summary.append(
            f"slide {idx:02d}: text_shapes={text_shapes}, "
            f"has_addition={has_addition}, has_flow={has_flow}, has_screenshot={has_screenshot}"
        )
    print("\n".join(summary))
    if issues:
        print("\nISSUES")
        print("\n".join(issues))
        return 1
    print("\nSTRUCTURE_QA_OK")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
